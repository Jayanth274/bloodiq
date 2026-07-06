import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Set slide dimensions to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette (Premium Dark Mode)
    BG_COLOR = RGBColor(15, 23, 42)      # Slate 900: #0F172A
    TEXT_WHITE = RGBColor(248, 250, 252) # Slate 50: #F8FAFC
    TEXT_MUTED = RGBColor(148, 163, 184) # Slate 400: #94A3B8
    ACCENT_RED = RGBColor(225, 29, 72)   # Rose 600: #E11D48
    ACCENT_BLUE = RGBColor(56, 189, 248)  # Sky 400: #38BDF8
    CARD_BG = RGBColor(30, 41, 59)       # Slate 800: #1E293B

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_title(slide, text, subtitle_text=None):
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(1.2))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_top = Inches(0)
        
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = 'Outfit'
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        
        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.name = 'Inter'
            p2.font.size = Pt(14)
            p2.font.color.rgb = ACCENT_BLUE
            p2.space_before = Pt(4)

    # ----------------- SLIDE 1: Title & Participant Details -----------------
    slide_layout = prs.slide_layouts[6] # Blank layout
    slide1 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide1)

    # Big Brand Header
    brand_box = slide1.shapes.add_textbox(Inches(0.75), Inches(1.2), Inches(11.83), Inches(2.2))
    tf1 = brand_box.text_frame
    tf1.word_wrap = True
    p_brand = tf1.paragraphs[0]
    p_brand.text = "BloodIQ"
    p_brand.font.name = 'Outfit'
    p_brand.font.size = Pt(64)
    p_brand.font.bold = True
    p_brand.font.color.rgb = ACCENT_RED

    p_tagline = tf1.add_paragraph()
    p_tagline.text = "GPU-Accelerated 72-Hour Blood Shortage Forecasting & Donor Mobilization Dashboard"
    p_tagline.font.name = 'Inter'
    p_tagline.font.size = Pt(18)
    p_tagline.font.color.rgb = TEXT_WHITE
    p_tagline.space_before = Pt(10)

    # Details Box (Card)
    details_shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(3.8), Inches(11.83), Inches(2.8))
    details_shape.fill.solid()
    details_shape.fill.fore_color.rgb = CARD_BG
    details_shape.line.color.rgb = ACCENT_BLUE
    details_shape.line.width = Pt(1.5)

    tf_det = details_shape.text_frame
    tf_det.word_wrap = True
    tf_det.margin_left = Inches(0.4)
    tf_det.margin_top = Inches(0.3)

    p_part = tf_det.paragraphs[0]
    p_part.text = "Participant Details"
    p_part.font.name = 'Outfit'
    p_part.font.size = Pt(20)
    p_part.font.bold = True
    p_part.font.color.rgb = ACCENT_BLUE

    p_name = tf_det.add_paragraph()
    p_name.text = "• Participant Name: Sri Sai Jayanth"
    p_name.font.name = 'Inter'
    p_name.font.size = Pt(15)
    p_name.font.color.rgb = TEXT_WHITE
    p_name.space_before = Pt(12)

    p_prob = tf_det.add_paragraph()
    p_prob.text = "• Problem Statement: Fatal emergency delays caused by unpredictable blood stock levels, delayed donor outreach, and zero real-time spatial predictive tracking across hospital networks."
    p_prob.font.name = 'Inter'
    p_prob.font.size = Pt(15)
    p_prob.font.color.rgb = TEXT_WHITE
    p_prob.space_before = Pt(12)

    # ----------------- SLIDE 2: Brief about the idea -----------------
    slide2 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide2)
    add_title(slide2, "Brief About the Idea", "Revolutionizing Emergency Blood Logistics")

    idea_box = slide2.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(11.83), Inches(4.5))
    tf2 = idea_box.text_frame
    tf2.word_wrap = True

    p = tf2.paragraphs[0]
    p.text = "What is BloodIQ?"
    p.font.name = 'Outfit'
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_RED

    p2 = tf2.add_paragraph()
    p2.text = "A modern predictive dashboard combining high-performance GPU data processing with conversational AI to eliminate blood supply blindspots."
    p2.font.name = 'Inter'
    p2.font.size = Pt(16)
    p2.font.color.rgb = TEXT_WHITE
    p2.space_before = Pt(8)

    p3 = tf2.add_paragraph()
    p3.text = "Key Value Proposition:"
    p3.font.name = 'Outfit'
    p3.font.size = Pt(22)
    p3.font.bold = True
    p3.font.color.rgb = ACCENT_RED
    p3.space_before = Pt(24)

    bullet_points = [
        "Proactive Shortage Prevention: Computes real-time Criticality Scores (0-100) and predicts stock depletion 72 hours ahead.",
        "Precision Donor Mobilization: Generates an emergency priority list matching recipient location, blood rarity, and donor recency.",
        "Zero-Learning NLP Chat: Non-technical staff query inventory state in plain English via a built-in Gemini Assistant."
    ]

    for pt in bullet_points:
        p_bullet = tf2.add_paragraph()
        p_bullet.text = "• " + pt
        p_bullet.font.name = 'Inter'
        p_bullet.font.size = Pt(16)
        p_bullet.font.color.rgb = TEXT_WHITE
        p_bullet.space_before = Pt(10)

    # ----------------- SLIDE 3: Solution Explanation -----------------
    slide3 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide3)
    add_title(slide3, "Translating Idea to Impact", "Core Solution Details")

    # Three-Column Layout for the three prompt questions
    col_width = Inches(3.6)
    col_gap = Inches(0.4)
    start_left = Inches(0.75)
    top_pos = Inches(2.2)
    height = Inches(4.5)

    columns_data = [
        {
            "title": "1. Technical Approach",
            "points": [
                "NVIDIA RAPIDS (cuDF): Processes 700MB+ inventory timeseries in seconds on T4 GPU.",
                "Google BigQuery: Central warehouse storing historical donation metrics & forecasts.",
                "Google Cloud Storage: Live automated ingestion bucket.",
                "Gemini API: Resolves natural language user search queries."
            ]
        },
        {
            "title": "2. Real-World Impact",
            "points": [
                "Reduces Stockouts: Flags critical stock depletion 3 days prior, allowing transfers.",
                "Speeds Mobilization: Targets rare O- donors in the exact postal code immediately.",
                "Saves Lives: Minimizes dispatch delay from hours to single-digit minutes."
            ]
        },
        {
            "title": "3. Data to Insight",
            "points": [
                "Ingestion: Uploads daily blood logs to cloud buckets.",
                "Scoring Engine: Calculates inventory status & penalties.",
                "Interactive Dashboard: Map UI represents danger zones (Red) and safety zones (Green) visually."
            ]
        }
    ]

    for i, col in enumerate(columns_data):
        left_pos = start_left + i * (col_width + col_gap)
        col_shape = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, top_pos, col_width, height)
        col_shape.fill.solid()
        col_shape.fill.fore_color.rgb = CARD_BG
        col_shape.line.color.rgb = ACCENT_BLUE if i == 0 else ACCENT_RED if i == 1 else TEXT_MUTED
        col_shape.line.width = Pt(1.5)

        tf_c = col_shape.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = Inches(0.2)
        tf_c.margin_right = Inches(0.2)
        tf_c.margin_top = Inches(0.2)

        p_t = tf_c.paragraphs[0]
        p_t.text = col["title"]
        p_t.font.name = 'Outfit'
        p_t.font.size = Pt(18)
        p_t.font.bold = True
        p_t.font.color.rgb = ACCENT_BLUE

        for pt in col["points"]:
            p_pt = tf_c.add_paragraph()
            p_pt.text = "• " + pt
            p_pt.font.name = 'Inter'
            p_pt.font.size = Pt(13)
            p_pt.font.color.rgb = TEXT_WHITE
            p_pt.space_before = Pt(8)

    # ----------------- SLIDE 4: Opportunities & USP -----------------
    slide4 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide4)
    add_title(slide4, "Market Opportunities & USP", "Standing Out From the Rest")

    # Left Column: Differentiator, Right Column: USP
    col_width_half = Inches(5.6)
    left_diff = Inches(0.75)
    left_usp = Inches(6.98)

    # Differentiator
    diff_shape = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_diff, top_pos, col_width_half, height)
    diff_shape.fill.solid()
    diff_shape.fill.fore_color.rgb = CARD_BG
    diff_shape.line.color.rgb = TEXT_MUTED
    diff_shape.line.width = Pt(1.5)

    tf_diff = diff_shape.text_frame
    tf_diff.word_wrap = True
    tf_diff.margin_left = Inches(0.3)
    tf_diff.margin_top = Inches(0.3)

    p_diff_t = tf_diff.paragraphs[0]
    p_diff_t.text = "How BloodIQ Differs from Existing Ideas"
    p_diff_t.font.name = 'Outfit'
    p_diff_t.font.size = Pt(20)
    p_diff_t.font.bold = True
    p_diff_t.font.color.rgb = ACCENT_BLUE

    diff_pts = [
        "Traditional: Passive directory listing static units.",
        "BloodIQ: Active forecasting engine showing 72h future depletion state.",
        "Traditional: Generic list of donors sorted by name or sign-up date.",
        "BloodIQ: Smart proximity & rarity priority lists generated dynamically.",
        "Traditional: Hard-to-use search UI with dropdown filters.",
        "BloodIQ: Gemini assistant resolves questions instantly like a colleague."
    ]
    for pt in diff_pts:
        p_pt = tf_diff.add_paragraph()
        p_pt.text = pt
        p_pt.font.name = 'Inter'
        p_pt.font.size = Pt(14)
        p_pt.font.color.rgb = TEXT_WHITE
        p_pt.space_before = Pt(10)

    # USP
    usp_shape = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_usp, top_pos, col_width_half, height)
    usp_shape.fill.solid()
    usp_shape.fill.fore_color.rgb = CARD_BG
    usp_shape.line.color.rgb = ACCENT_RED
    usp_shape.line.width = Pt(1.5)

    tf_usp = usp_shape.text_frame
    tf_usp.word_wrap = True
    tf_usp.margin_left = Inches(0.3)
    tf_usp.margin_top = Inches(0.3)

    p_usp_t = tf_usp.paragraphs[0]
    p_usp_t.text = "Unique Selling Propositions (USPs)"
    p_usp_t.font.name = 'Outfit'
    p_usp_t.font.size = Pt(20)
    p_usp_t.font.bold = True
    p_usp_t.font.color.rgb = ACCENT_RED

    usp_pts = [
        "GPU-Accelerated Analytics: cuDF processes heavy timeseries and calculates donor rankings under 3 seconds (10x faster than Pandas).",
        "Predictive Severity Index: Dynamically adjusts penalties based on days left for local donor camps.",
        "Hybrid Gemini AI Interface: Allows voice-like natural queries mapping to underlying Structured Data API responses.",
        "Single-Screen Situational Awareness: Fully unified React map UI with sidebar feeds."
    ]
    for pt in usp_pts:
        p_pt = tf_usp.add_paragraph()
        p_pt.text = "✔ " + pt
        p_pt.font.name = 'Inter'
        p_pt.font.size = Pt(14)
        p_pt.font.color.rgb = TEXT_WHITE
        p_pt.space_before = Pt(10)

    # ----------------- SLIDE 5: List of features offered -----------------
    slide5 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide5)
    add_title(slide5, "Features Offered by the Solution", "Ready for Real-world Implementation")

    feat_box = slide5.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(11.83), Inches(4.5))
    tf5 = feat_box.text_frame
    tf5.word_wrap = True

    features = [
        ("72-Hour Predictive Alerts", "Automatically triggers warning notifications on the dashboard 72 hours before stock hits absolute zero."),
        ("Interactive Leaflet Map Dashboard", "Visualizes blood availability across cities using GPS coordinates. Green = Safe, Yellow = Low, Red = Critical shortage."),
        ("Ranked Donor Mobilization Feed", "Provides immediate emergency contact details, prioritizing donors based on travel distance, blood type rarity, and historical donation intervals."),
        ("Built-in Conversational Assistant", "Resolves complex operational questions (e.g., 'Which O- banks are critical in Mumbai?') in seconds using LLM processing."),
        ("Unified GPU Pipeline Benchmarks", "Uses cuDF (NVIDIA RAPIDS) to run the clean-and-score analytics pipeline locally, reducing query latency and ensuring instant sync.")
    ]

    for title, desc in features:
        p_title = tf5.add_paragraph()
        p_title.text = "★ " + title
        p_title.font.name = 'Outfit'
        p_title.font.size = Pt(18)
        p_title.font.bold = True
        p_title.font.color.rgb = ACCENT_BLUE
        if tf5.paragraphs[0] != p_title:
            p_title.space_before = Pt(12)

        p_desc = tf5.add_paragraph()
        p_desc.text = "   " + desc
        p_desc.font.name = 'Inter'
        p_desc.font.size = Pt(14)
        p_desc.font.color.rgb = TEXT_WHITE
        p_desc.space_before = Pt(2)

    # ----------------- SLIDE 6: Process Flow Diagram -----------------
    slide6 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide6)
    add_title(slide6, "Process Flow Diagram", "Data Lifecycle & Decision Pipeline")

    # Creating flowchart boxes
    box_width = Inches(2.0)
    box_height = Inches(1.5)
    y_pos = Inches(3.0)

    steps = [
        {"name": "1. Ingestion", "desc": "Raw CSV records (Banks, Donors, Donations) loaded to GCS."},
        {"name": "2. GPU Pipeline", "desc": "cuDF processes, joins datasets & computes metrics."},
        {"name": "3. ML & Scoring", "desc": "Generates 72h forecasts and donor priorities."},
        {"name": "4. API Layer", "desc": "Express REST API serves processed JSON responses."},
        {"name": "5. UI Layer", "desc": "React Map displays alerts & routes query requests."}
    ]

    for idx, step in enumerate(steps):
        x_pos = Inches(0.75) + idx * Inches(2.4)
        # Box
        shape = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, y_pos, box_width, box_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = ACCENT_RED
        shape.line.width = Pt(2)

        tf_s = shape.text_frame
        tf_s.word_wrap = True
        tf_s.margin_left = Inches(0.1)
        tf_s.margin_right = Inches(0.1)
        tf_s.margin_top = Inches(0.1)

        p_n = tf_s.paragraphs[0]
        p_n.text = step["name"]
        p_n.font.name = 'Outfit'
        p_n.font.size = Pt(14)
        p_n.font.bold = True
        p_n.font.color.rgb = ACCENT_BLUE
        p_n.alignment = PP_ALIGN.CENTER

        p_d = tf_s.add_paragraph()
        p_d.text = step["desc"]
        p_d.font.name = 'Inter'
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = TEXT_WHITE
        p_d.space_before = Pt(4)
        p_d.alignment = PP_ALIGN.CENTER

        # Connector Arrow
        if idx < len(steps) - 1:
            arrow = slide6.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x_pos + box_width + Inches(0.1), y_pos + Inches(0.5), Inches(0.2), Inches(0.5))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT_BLUE
            arrow.line.fill.background()

    # ----------------- SLIDE 7: Wireframes/Mock diagrams -----------------
    slide7 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide7)
    add_title(slide7, "Dashboard Interface Wireframe Layout", "Optimized for Single-Screen Crisis Resolution")

    mock_shape = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(2.2), Inches(11.83), Inches(4.5))
    mock_shape.fill.solid()
    mock_shape.fill.fore_color.rgb = CARD_BG
    mock_shape.line.color.rgb = ACCENT_BLUE
    mock_shape.line.width = Pt(1.5)

    tf_mock = mock_shape.text_frame
    tf_mock.word_wrap = True
    tf_mock.margin_left = Inches(0.4)
    tf_mock.margin_top = Inches(0.3)

    p_mock_t = tf_mock.paragraphs[0]
    p_mock_t.text = "Dashboard UI Layout Architecture"
    p_mock_t.font.name = 'Outfit'
    p_mock_t.font.size = Pt(20)
    p_mock_t.font.bold = True
    p_mock_t.font.color.rgb = ACCENT_RED

    mock_pts = [
        "1. Left Sidebar: Crisis Alert Feed showing urgent depleted banks and expected stockout times.",
        "2. Central Content Area: Interactive Leaflet.js map with color-coded circular pins indicating stock severity. Filter controls for city and blood type.",
        "3. Right Drawer: Opens upon pin click. Shows selected bank's details, nearest backup inventory sources, and a conversational ChatBox.",
        "4. Conversational ChatBox: Placed in the panel, lets officers ask Gemini questions like 'Find A- donors close to City Hospital' and see immediate matches."
    ]
    for pt in mock_pts:
        p_pt = tf_mock.add_paragraph()
        p_pt.text = pt
        p_pt.font.name = 'Inter'
        p_pt.font.size = Pt(15)
        p_pt.font.color.rgb = TEXT_WHITE
        p_pt.space_before = Pt(12)

    # ----------------- SLIDE 8: Architecture Diagram -----------------
    slide8 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide8)
    add_title(slide8, "Solution Architecture Diagram", "Robust, Serverless, and Scalable Stack")

    # Drawing components:
    # GCS Bucket -> BigQuery Dataset -> cuDF GPU Pipeline -> Node.js API -> React Frontend + Gemini API
    comp_w = Inches(2.2)
    comp_h = Inches(1.2)

    # Row 1: GCS Bucket (Ingestion)
    s_gcs = slide8.shapes.add_shape(MSO_SHAPE.CAN, Inches(0.75), Inches(3.0), comp_w, comp_h)
    s_gcs.fill.solid()
    s_gcs.fill.fore_color.rgb = CARD_BG
    s_gcs.line.color.rgb = ACCENT_BLUE
    s_gcs.text_frame.text = "Google Cloud Storage\n(bloodiq-raw-data)"
    s_gcs.text_frame.paragraphs[0].font.size = Pt(11)
    s_gcs.text_frame.paragraphs[0].font.bold = True

    # Arrow to BigQuery
    arrow1 = slide8.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.05), Inches(3.35), Inches(0.2), Inches(0.5))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = TEXT_MUTED

    # BigQuery
    s_bq = slide8.shapes.add_shape(MSO_SHAPE.CAN, Inches(3.35), Inches(3.0), comp_w, comp_h)
    s_bq.fill.solid()
    s_bq.fill.fore_color.rgb = CARD_BG
    s_bq.line.color.rgb = ACCENT_BLUE
    s_bq.text_frame.text = "Google BigQuery\n(bloodiq_data)"
    s_bq.text_frame.paragraphs[0].font.size = Pt(11)
    s_bq.text_frame.paragraphs[0].font.bold = True

    # Arrow to cuDF
    arrow2 = slide8.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.65), Inches(3.35), Inches(0.2), Inches(0.5))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = TEXT_MUTED

    # cuDF RAPIDS
    s_cudf = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.95), Inches(3.0), comp_w, comp_h)
    s_cudf.fill.solid()
    s_cudf.fill.fore_color.rgb = CARD_BG
    s_cudf.line.color.rgb = ACCENT_RED
    s_cudf.text_frame.text = "NVIDIA cuDF GPU Pipeline\n(72h Forecast & Donor Rank)"
    s_cudf.text_frame.paragraphs[0].font.size = Pt(11)
    s_cudf.text_frame.paragraphs[0].font.bold = True

    # Arrow to Node Express
    arrow3 = slide8.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.25), Inches(3.35), Inches(0.2), Inches(0.5))
    arrow3.fill.solid()
    arrow3.fill.fore_color.rgb = TEXT_MUTED

    # Node.js
    s_node = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.55), Inches(3.0), comp_w, comp_h)
    s_node.fill.solid()
    s_node.fill.fore_color.rgb = CARD_BG
    s_node.line.color.rgb = ACCENT_BLUE
    s_node.text_frame.text = "Node.js / Express API\n(4 core endpoints)"
    s_node.text_frame.paragraphs[0].font.size = Pt(11)
    s_node.text_frame.paragraphs[0].font.bold = True

    # Arrow to React
    arrow4 = slide8.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(10.85), Inches(3.35), Inches(0.2), Inches(0.5))
    arrow4.fill.solid()
    arrow4.fill.fore_color.rgb = TEXT_MUTED

    # React Frontend
    s_react = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.15), Inches(3.0), comp_w, comp_h)
    s_react.fill.solid()
    s_react.fill.fore_color.rgb = CARD_BG
    s_react.line.color.rgb = ACCENT_RED
    s_react.text_frame.text = "React 18 Dashboard\n(Leaflet Map + Sidebar)"
    s_react.text_frame.paragraphs[0].font.size = Pt(11)
    s_react.text_frame.paragraphs[0].font.bold = True

    # Gemini Assistant (Double Arrow connecting with Node API)
    s_gemini = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.55), Inches(5.0), comp_w, comp_h)
    s_gemini.fill.solid()
    s_gemini.fill.fore_color.rgb = CARD_BG
    s_gemini.line.color.rgb = ACCENT_RED
    s_gemini.text_frame.text = "Google Gemini API\n(NL Query Engine)"
    s_gemini.text_frame.paragraphs[0].font.size = Pt(11)
    s_gemini.text_frame.paragraphs[0].font.bold = True

    arrow5 = slide8.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(9.55), Inches(4.3), Inches(0.2), Inches(0.65))
    arrow5.fill.solid()
    arrow5.fill.fore_color.rgb = TEXT_MUTED

    # ----------------- SLIDE 9: Technologies Used & Scalability -----------------
    slide9 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide9)
    add_title(slide9, "Technologies & Services Integration", "Why this Stack Matters for Scalability")

    tech_box = slide9.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(11.83), Inches(4.5))
    tf9 = tech_box.text_frame
    tf9.word_wrap = True

    techs = [
        ("NVIDIA cuDF & RAPIDS", "GPU acceleration allows running timeseries joins and aggregations on millions of blood bank records in <3s. Resolves performance limits of pandas, enabling massive scale during regional health emergencies."),
        ("Google BigQuery", "Highly scalable analytics dataset engine supporting structural SQL storage for fast recovery, data auditing, and clean synchronization of historical metrics."),
        ("Google Cloud Storage (GCS)", "Acts as the ingestion landing zone. Fast, durable, and easily integrated with Google Cloud Functions to trigger analytical pipelines automatically on upload."),
        ("Google Gemini API", "Enables intuitive voice/text querying, parsing messy phrases (e.g. 'need A+ now near Pune') into clean structured API queries without database code lookup.")
    ]

    for title, desc in techs:
        p_title = tf9.add_paragraph()
        p_title.text = "■ " + title
        p_title.font.name = 'Outfit'
        p_title.font.size = Pt(17)
        p_title.font.bold = True
        p_title.font.color.rgb = ACCENT_RED
        if tf9.paragraphs[0] != p_title:
            p_title.space_before = Pt(10)

        p_desc = tf9.add_paragraph()
        p_desc.text = "   " + desc
        p_desc.font.name = 'Inter'
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = TEXT_WHITE
        p_desc.space_before = Pt(2)

    # ----------------- SLIDE 10: Prototype Snapshots -----------------
    slide10 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide10)
    add_title(slide10, "Prototype Snapshots & Verified Backend Endpoints", "Evidence of Accelerated Delivery")

    snap_box = slide10.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(11.83), Inches(4.5))
    tf10 = snap_box.text_frame
    tf10.word_wrap = True

    p_sub = tf10.paragraphs[0]
    p_sub.text = "Clean REST APIs supporting the UI (all verified to work on localhost:3001):"
    p_sub.font.name = 'Outfit'
    p_sub.font.size = Pt(18)
    p_sub.font.bold = True
    p_sub.font.color.rgb = ACCENT_BLUE

    endpoints = [
        ("GET /api/availability", "Returns current inventory states grouped by cities. Integrated with Leaflet visual mapping pins."),
        ("GET /api/forecast", "Outputs 72-hour forecast severity lists, calculated in Python cuDF, sorted by highest criticality first."),
        ("GET /api/donors", "Retrieves active donors, ranked using dynamic distance, donation intervals, and blood rarity weights."),
        ("POST /api/query", "Accepts plain-text questions, routes to Gemini Assistant, and returns clear, conversational responses.")
    ]

    for ep, desc in endpoints:
        p_ep = tf10.add_paragraph()
        p_ep.text = "● " + ep
        p_ep.font.name = 'Outfit'
        p_ep.font.size = Pt(15)
        p_ep.font.bold = True
        p_ep.font.color.rgb = ACCENT_RED
        p_ep.space_before = Pt(8)

        p_ed = tf10.add_paragraph()
        p_ed.text = "   " + desc
        p_ed.font.name = 'Inter'
        p_ed.font.size = Pt(13)
        p_ed.font.color.rgb = TEXT_WHITE
        p_ed.space_before = Pt(2)

    # ----------------- SLIDE 11: Thank you -----------------
    slide11 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide11)

    thank_box = slide11.shapes.add_textbox(Inches(0.75), Inches(2.5), Inches(11.83), Inches(3.0))
    tf11 = thank_box.text_frame
    tf11.word_wrap = True

    p_t1 = tf11.paragraphs[0]
    p_t1.text = "Thank You!"
    p_t1.font.name = 'Outfit'
    p_t1.font.size = Pt(64)
    p_t1.font.bold = True
    p_t1.font.color.rgb = ACCENT_RED
    p_t1.alignment = PP_ALIGN.CENTER

    p_t2 = tf11.add_paragraph()
    p_t2.text = "Proactive emergency response powered by NVIDIA & Google Cloud"
    p_t2.font.name = 'Inter'
    p_t2.font.size = Pt(20)
    p_t2.font.color.rgb = TEXT_WHITE
    p_t2.space_before = Pt(10)
    p_t2.alignment = PP_ALIGN.CENTER

    # Save presentation to root directory relative path
    prs.save('BloodIQ_Submission.pptx')
    print("Presentation generated successfully at BloodIQ_Submission.pptx")

if __name__ == "__main__":
    create_presentation()
